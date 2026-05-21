#!/usr/bin/env python3
"""Create an editable PowerPoint version of the Fig. 2-1 VQ-VAE flow diagram."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/home/llb/HunyuanWorld-Voyager/bishe")
OUT_DIR = ROOT / "carla_project/outputs/figures/thesis_visuals_20260519_compact"
OUT_PATH = OUT_DIR / "fig2_1_vqvae_flow_editable.pptx"

FONT = "Microsoft YaHei"
BLUE = RGBColor(62, 102, 153)
BLUE_LIGHT = RGBColor(226, 237, 250)
BLUE_MID = RGBColor(156, 183, 222)
GREEN = RGBColor(58, 135, 66)
GREEN_LIGHT = RGBColor(232, 245, 229)
GREEN_MID = RGBColor(174, 218, 164)
ORANGE = RGBColor(198, 109, 48)
ORANGE_LIGHT = RGBColor(252, 235, 219)
GRAY = RGBColor(96, 96, 96)
WHITE = RGBColor(255, 255, 255)


def set_shape_text(shape, text: str, size: float = 12, color=RGBColor(0, 0, 0), bold: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold


def add_text(slide, x, y, w, h, text, size=12, color=RGBColor(0, 0, 0), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_text(box, text, size=size, color=color, bold=bold)
    return box


def add_number(slide, num: str, x: float, y: float, color=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.28), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    set_shape_text(shape, num, size=11, color=WHITE, bold=True)


def add_round_rect(slide, x, y, w, h, fill, line, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.2)
    return shp


def add_right_arrow(slide, x, y, w=0.36, h=0.22, color=BLUE) -> None:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.color.rgb = color


def add_image_icon(slide, x, y, w=1.15, h=1.15) -> None:
    box = add_round_rect(slide, x, y, w, h, BLUE_LIGHT, BLUE)
    # sky, ground, mountains and sun are all editable shapes.
    ground = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.03), Inches(y + h * 0.68), Inches(w - 0.06), Inches(h * 0.27))
    ground.fill.solid()
    ground.fill.fore_color.rgb = RGBColor(135, 183, 222)
    ground.line.color.rgb = RGBColor(135, 183, 222)
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.14), Inches(y + 0.18), Inches(0.22), Inches(0.22))
    sun.fill.solid()
    sun.fill.fore_color.rgb = RGBColor(250, 197, 67)
    sun.line.color.rgb = RGBColor(230, 170, 40)
    m1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 0.18), Inches(y + 0.43), Inches(0.55), Inches(0.46))
    m1.fill.solid()
    m1.fill.fore_color.rgb = RGBColor(123, 184, 104)
    m1.line.color.rgb = GREEN
    m2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 0.53), Inches(y + 0.34), Inches(0.58), Inches(0.56))
    m2.fill.solid()
    m2.fill.fore_color.rgb = RGBColor(105, 169, 91)
    m2.line.color.rgb = GREEN
    box.line.width = Pt(1.3)


def add_encoder_icon(slide, x, y, mirror=False, w=1.05, h=1.15) -> None:
    add_round_rect(slide, x, y, w, h, RGBColor(248, 251, 255), BLUE)
    widths = [0.15, 0.13, 0.11, 0.09]
    heights = [0.72, 0.55, 0.40, 0.26]
    if mirror:
        xs = [x + 0.20, x + 0.39, x + 0.58, x + 0.77]
        widths = widths[::-1]
        heights = heights[::-1]
    else:
        xs = [x + 0.20, x + 0.41, x + 0.60, x + 0.77]
    for xi, wi, hi in zip(xs, widths, heights):
        bar = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(xi), Inches(y + (h - hi) / 2), Inches(wi), Inches(hi))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE_MID
        bar.line.color.rgb = BLUE


def add_feature_grid(slide, x, y, w=1.12, h=1.12, line=BLUE, fill=RGBColor(244, 249, 255)) -> None:
    add_round_rect(slide, x, y, w, h, fill, line, radius=False)
    rows = 6
    cols = 6
    for i in range(1, cols):
        lx = x + w * i / cols
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx), Inches(y), Inches(0.005), Inches(h))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = line
        line_shape.line.color.rgb = line
    for i in range(1, rows):
        ly = y + h * i / rows
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(ly), Inches(w), Inches(0.005))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = line
        line_shape.line.color.rgb = line


def add_quant_icon(slide, x, y, w=1.15, h=1.15) -> None:
    add_round_rect(slide, x, y, w, h, GREEN_LIGHT, GREEN)
    left_points = [(x + 0.25, y + 0.34), (x + 0.38, y + 0.52), (x + 0.22, y + 0.68), (x + 0.43, y + 0.78), (x + 0.33, y + 0.88)]
    right_points = [(x + 0.78, y + 0.35), (x + 0.91, y + 0.49), (x + 0.70, y + 0.66), (x + 0.88, y + 0.78), (x + 0.76, y + 0.89)]
    for px, py in left_points + right_points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(0.07), Inches(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.color.rgb = GREEN
    add_right_arrow(slide, x + 0.48, y + 0.59, 0.24, 0.13, GREEN)


def add_token_table(slide, x, y, w=1.52, h=1.42) -> None:
    rows, cols = 5, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    values = [
        ["102", "7", "345", "…", "88"],
        ["12", "2048", "3", "…", "77"],
        ["409", "1", "256", "…", "30"],
        ["⋮", "⋮", "⋮", "⋱", "⋮"],
        ["66", "199", "512", "…", "8"],
    ]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = values[r][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN_LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT
            p.font.size = Pt(8.5)
            p.font.color.rgb = RGBColor(22, 92, 31)
            p.font.bold = r != 3


def add_codebook_table(slide, x, y, w=2.10, h=1.23) -> None:
    rows, cols = 5, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    values = [
        ["0", "e₀¹", "e₀²", "…", "e₀ᵈ"],
        ["1", "e₁¹", "e₁²", "…", "e₁ᵈ"],
        ["2", "e₂¹", "e₂²", "…", "e₂ᵈ"],
        ["⋮", "⋮", "⋮", "⋱", "⋮"],
        ["4095", "e₄₀₉₅¹", "e₄₀₉₅²", "…", "e₄₀₉₅ᵈ"],
    ]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = values[r][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN_LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT
            p.font.size = Pt(7.8)
            p.font.color.rgb = RGBColor(22, 92, 31)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Module positions.
    y_icon = 2.02
    xs = {
        "input": 0.25,
        "enc": 1.80,
        "latent": 3.23,
        "quant": 4.82,
        "token": 6.45,
        "vec": 8.25,
        "dec": 10.02,
        "out": 11.55,
    }

    # Title-like labels and number markers.
    labels = [
        ("1", "输入图像 x\n256×256×3", xs["input"] + 0.43, BLUE),
        ("2", "编码器\nEncoder", xs["enc"] + 0.39, BLUE),
        ("3", "连续潜在特征 zₑ(x)\n32×32×d", xs["latent"] + 0.40, BLUE),
        ("4", "向量量化\nQuantization", xs["quant"] + 0.43, GREEN),
        ("6", "离散 token 网格 z\n32×32", xs["token"] + 0.55, GREEN),
        ("7", "码本向量 z_q(x)\n32×32×d", xs["vec"] + 0.45, GREEN),
        ("8", "解码器\nDecoder", xs["dec"] + 0.38, BLUE),
        ("9", "重建图像 x̂\n256×256×3", xs["out"] + 0.42, BLUE),
    ]
    for num, label, cx, color in labels:
        add_number(slide, num, cx, 0.45, color)
        add_text(slide, cx - 0.55, 0.82, 1.38, 0.55, label, size=10.2, color=RGBColor(0, 0, 0), bold=True)

    # Main flow.
    add_image_icon(slide, xs["input"], y_icon)
    add_encoder_icon(slide, xs["enc"], y_icon)
    add_feature_grid(slide, xs["latent"], y_icon + 0.02, w=1.12, h=1.12, line=BLUE)
    add_quant_icon(slide, xs["quant"], y_icon)
    add_token_table(slide, xs["token"], y_icon - 0.08, w=1.52, h=1.40)
    add_feature_grid(slide, xs["vec"], y_icon + 0.02, w=1.12, h=1.12, line=GREEN, fill=GREEN_LIGHT)
    add_encoder_icon(slide, xs["dec"], y_icon, mirror=True)
    add_image_icon(slide, xs["out"], y_icon)

    arrow_y = y_icon + 0.47
    add_right_arrow(slide, xs["input"] + 1.20, arrow_y)
    add_right_arrow(slide, xs["enc"] + 1.12, arrow_y)
    add_right_arrow(slide, xs["latent"] + 1.18, arrow_y)
    add_right_arrow(slide, xs["quant"] + 1.20, arrow_y)
    add_right_arrow(slide, xs["token"] + 1.58, arrow_y, color=GREEN)
    add_right_arrow(slide, xs["vec"] + 1.18, arrow_y)
    add_right_arrow(slide, xs["dec"] + 1.12, arrow_y)

    # Codebook branch.
    add_number(slide, "5", 4.72, 4.45, GREEN)
    add_text(slide, 5.02, 4.42, 1.65, 0.42, "码本 Codebook E\n4096×d", size=10.0, color=RGBColor(22, 92, 31), bold=True)
    add_codebook_table(slide, 4.88, 4.92, w=2.20, h=1.28)
    up_arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(5.63), Inches(3.72), Inches(0.20), Inches(0.72))
    up_arrow.fill.solid()
    up_arrow.fill.fore_color.rgb = GREEN
    up_arrow.line.color.rgb = GREEN
    add_text(slide, 5.82, 3.85, 0.92, 0.30, "最近邻查找", size=9.3, color=GREEN, bold=True)
    add_text(slide, 6.82, 5.05, 0.95, 0.58, "EMA 更新\n死码重置", size=8.0, color=GREEN)

    # Small footer note keeps the figure tied to the thesis setting.
    add_text(
        slide,
        0.34,
        6.92,
        12.5,
        0.28,
        "图中数值对应本文 f=8 设置：256×256 图像编码为 32×32 token 网格，码本容量为 4096。",
        size=9.0,
        color=GRAY,
    )

    prs.save(OUT_PATH)
    print(f"saved {OUT_PATH}")


if __name__ == "__main__":
    main()
